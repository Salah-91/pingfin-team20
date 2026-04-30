"""
PingFin Team 20 — 25-slide presentation generator
Run: python docs/generate_presentation.py
Output: docs/PingFin-Team20-Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ── Color palette: Banking — deep navy + gold accent ───────────────────
NAVY        = RGBColor(0x0F, 0x1B, 0x3D)   # primary background
NAVY_LIGHT  = RGBColor(0x1E, 0x2F, 0x5C)   # cards on dark
ICE_BLUE    = RGBColor(0xCA, 0xDC, 0xFC)   # text on dark / accents
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT  = RGBColor(0xE2, 0xEC, 0xFF)
TEXT_MUTED  = RGBColor(0x7C, 0x93, 0xB8)
GOLD        = RGBColor(0xF2, 0xC9, 0x4C)   # accent — big numbers, CTA
GREEN_OK    = RGBColor(0x22, 0xD3, 0xA0)
RED_ERR     = RGBColor(0xF2, 0x57, 0x57)
AMBER       = RGBColor(0xF5, 0xA6, 0x23)
DARK_BG     = RGBColor(0x07, 0x0B, 0x14)   # extra dark for title slides

HEADER_FONT = "Calibri"
BODY_FONT   = "Calibri"
MONO_FONT   = "Consolas"


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ──────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()

def stroke(shape, rgb, width=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width)

def bg(slide, rgb):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    fill(bg, rgb)
    bg.shadow.inherit = False
    return bg

def text_box(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False,
             color=TEXT_LIGHT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, list):
        for i, segment in enumerate(text):
            run = p.add_run() if i > 0 else p.add_run()
            if isinstance(segment, str):
                run.text = segment
                run.font.name = font
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = color
            else:
                run.text = segment["text"]
                run.font.name = segment.get("font", font)
                run.font.size = Pt(segment.get("size", size))
                run.font.bold = segment.get("bold", bold)
                run.font.color.rgb = segment.get("color", color)
    else:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tx

def add_lines(slide, x, y, w, h, lines, *, font=BODY_FONT, size=14,
              color=TEXT_LIGHT, bullet=False, line_space=1.15):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        if isinstance(line, str):
            run = p.add_run()
            run.text = ("•  " + line) if bullet else line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
        else:  # dict with options
            run = p.add_run()
            run.text = ("•  " + line["text"]) if bullet else line["text"]
            run.font.name = line.get("font", font)
            run.font.size = Pt(line.get("size", size))
            run.font.bold = line.get("bold", False)
            run.font.color.rgb = line.get("color", color)
    return tx

def card(slide, x, y, w, h, fill_rgb=NAVY_LIGHT, accent=None):
    """Card with optional left-accent bar."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(box, fill_rgb)
    box.shadow.inherit = False
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
        fill(bar, accent)
        bar.shadow.inherit = False
    return box

def chip(slide, x, y, text, color=GREEN_OK, text_color=NAVY):
    """Small pill/chip."""
    w = max(0.4, len(text) * 0.085 + 0.3)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32))
    fill(rect, color)
    rect.shadow.inherit = False
    rect.adjustments[0] = 0.5
    text_box(slide, x, y, w, 0.32, text, font=BODY_FONT, size=10, bold=True,
             color=text_color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

def page_header(slide, title, subtitle=None, num=None):
    """Top bar with title + slide-number indicator."""
    # Background
    bg(slide, NAVY)
    # Top accent line (gold, thin)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.12))
    fill(line, GOLD)
    line.shadow.inherit = False
    # Title
    text_box(slide, 0.6, 0.3, 11, 0.7, title, font=HEADER_FONT, size=28, bold=True, color=WHITE)
    if subtitle:
        text_box(slide, 0.6, 0.95, 11, 0.4, subtitle, font=BODY_FONT, size=13, color=TEXT_MUTED)
    # Slide number
    if num is not None:
        text_box(slide, 12.5, 7.0, 0.7, 0.3, f"{num} / 25", font=MONO_FONT, size=10,
                 color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    # Footer (small)
    text_box(slide, 0.6, 7.0, 8, 0.3, "PingFin · Team 20 · CEKVBE88 + HOMNBEB1",
             font=BODY_FONT, size=9, color=TEXT_MUTED)


def big_number(slide, x, y, value, label, color=GOLD, num_size=72, lbl_size=14):
    text_box(slide, x, y, 3, 1.4, value, font=HEADER_FONT, size=num_size, bold=True,
             color=color, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
    text_box(slide, x, y + 1.4, 3, 0.5, label, font=BODY_FONT, size=lbl_size,
             color=TEXT_MUTED, align=PP_ALIGN.LEFT)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
bg(s, DARK_BG)

# Decorative gold band on left
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SH)
fill(band, GOLD)
band.shadow.inherit = False

# Bank icon block
icon_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0))
fill(icon_box, GOLD)
icon_box.shadow.inherit = False
text_box(s, 1.0, 1.0, 1.0, 1.0, "🏦", font=HEADER_FONT, size=44, color=NAVY,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

text_box(s, 2.4, 1.05, 9, 0.5, "PINGFIN", font=HEADER_FONT, size=18, bold=True,
         color=GOLD)
text_box(s, 2.4, 1.55, 11, 1.4, "SEPA Payment System Simulation",
         font=HEADER_FONT, size=54, bold=True, color=WHITE)

text_box(s, 2.4, 3.1, 11, 0.5, "Team 20 · Workshop Software Engineering 2026",
         font=BODY_FONT, size=18, color=TEXT_LIGHT)

# Bank chips
chip(s, 2.4, 3.9, "BANK 1 · CEKVBE88 · port 8089", GOLD, NAVY)
chip(s, 5.6, 3.9, "BANK 2 · HOMNBEB1 · port 8090", GOLD, NAVY)

# Bottom info
text_box(s, 1.0, 6.4, 11, 0.4,
         "Salaheddine Sennouni  ·  Abdallah Azouagh  ·  Ayoub Abdeddoun  ·  Marwan Saidi",
         font=BODY_FONT, size=14, color=TEXT_MUTED)
text_box(s, 1.0, 6.8, 11, 0.3, "Coach: Rogier van der Linde · Polina Kozlova · Odisee",
         font=BODY_FONT, size=11, color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 2 — Team
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Het Team", "Wie heeft wat gedaan?", 2)

team = [
    ("SS", "Salaheddine Sennouni", "Backend & Architecture",
     "API endpoints · DB design · Background jobs · Security · Deployment"),
    ("AA", "Abdallah Azouagh", "Team Lead & Integration",
     "Project management · Trello · Cross-team testing · Verslag"),
    ("AB", "Ayoub Abdeddoun", "Validation & Database",
     "Validation rules · Foutcode-mapping · DB-schema · Test suite"),
    ("MS", "Marwan Saidi", "Frontend & Design",
     "GUI HTML/CSS · Live notifications · Auto-poll · Mockups"),
]

for i, (initials, name, role, what) in enumerate(team):
    col = i % 2
    row = i // 2
    cx = 0.6 + col * 6.3
    cy = 1.7 + row * 2.6

    card(s, cx, cy, 6.0, 2.3, NAVY_LIGHT, accent=GOLD)
    # Avatar circle
    av = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.4), Inches(cy + 0.4), Inches(1.0), Inches(1.0))
    fill(av, GOLD)
    av.shadow.inherit = False
    text_box(s, cx + 0.4, cy + 0.4, 1.0, 1.0, initials, font=HEADER_FONT, size=24,
             bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    text_box(s, cx + 1.6, cy + 0.4, 4.2, 0.4, name, font=HEADER_FONT, size=18,
             bold=True, color=WHITE)
    text_box(s, cx + 1.6, cy + 0.8, 4.2, 0.35, role, font=BODY_FONT, size=12,
             bold=True, color=GOLD)
    text_box(s, cx + 0.4, cy + 1.55, 5.4, 0.7, what, font=BODY_FONT, size=11,
             color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 3 — Probleemstelling
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Probleemstelling", "Wat is PingFin?", 3)

card(s, 0.6, 1.7, 12.1, 1.6, NAVY_LIGHT, accent=GOLD)
text_box(s, 0.9, 1.85, 11.5, 0.4, "Wat is SEPA?", font=HEADER_FONT, size=18,
         bold=True, color=GOLD)
text_box(s, 0.9, 2.3, 11.5, 1.0,
         "Single Euro Payments Area — een simulatie van het Europees betalingsverkeer waarbij banken "
         "PO's (Payment Orders) versturen via een centrale Clearing Bank. Elke transactie reist als "
         "een 'ping': OB → CB → BB → ACK terug.",
         font=BODY_FONT, size=14, color=TEXT_LIGHT)

# Three role cards
roles = [
    ("OB", "Originating Bank",
     "Verstuurt PO's namens de OA-rekeninghouder", GOLD),
    ("CB", "Clearing Bank",
     "Routeert berichten in het midden, valideert", ICE_BLUE),
    ("BB", "Beneficiary Bank",
     "Ontvangt PO, crediteert BA-rekening, stuurt ACK", GREEN_OK),
]
for i, (abbr, name, desc, accent) in enumerate(roles):
    cx = 0.6 + i * 4.05
    card(s, cx, 3.6, 3.85, 2.3, NAVY_LIGHT, accent=accent)
    text_box(s, cx + 0.3, 3.8, 1.0, 0.6, abbr, font=HEADER_FONT, size=42,
             bold=True, color=accent)
    text_box(s, cx + 0.3, 4.55, 3.4, 0.4, name, font=HEADER_FONT, size=16,
             bold=True, color=WHITE)
    text_box(s, cx + 0.3, 4.95, 3.4, 0.9, desc, font=BODY_FONT, size=12,
             color=TEXT_MUTED)

text_box(s, 0.6, 6.2, 12.1, 0.5,
         "Onze rol: 2 gewone banken die zowel OB als BB kunnen zijn",
         font=HEADER_FONT, size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 4 — Doelstellingen
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Doelstellingen", "Wat moest ons systeem kunnen?", 4)

goals = [
    ("📨", "PO's verzenden", "Genereer + valideer + verstuur via Clearing Bank"),
    ("📥", "PO's ontvangen", "Verwerk inkomende PO's + verstuur ACK terug"),
    ("⚖️", "Validatie", "BIC, IBAN, bedrag (≤€500), saldo, format-checks"),
    ("🛡️", "Foutafhandeling", "10 foutcodes (4001-4102) + timeout na 1u + auto-refund"),
    ("💻", "GUI", "Realtime monitoring + manuele PO + live notifications"),
    ("🔒", "Security", "Bearer-tokens, atomische transacties, SQL-veilig, XSS-veilig"),
]
for i, (icon, title, desc) in enumerate(goals):
    col = i % 3
    row = i // 3
    cx = 0.6 + col * 4.18
    cy = 1.8 + row * 2.5
    card(s, cx, cy, 4.0, 2.2, NAVY_LIGHT, accent=GOLD)
    text_box(s, cx + 0.3, cy + 0.25, 1.0, 0.7, icon, font=HEADER_FONT, size=32,
             color=GOLD, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, cx + 0.3, cy + 1.0, 3.6, 0.4, title, font=HEADER_FONT, size=17,
             bold=True, color=WHITE)
    text_box(s, cx + 0.3, cy + 1.45, 3.6, 0.7, desc, font=BODY_FONT, size=12,
             color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 5 — Architectuur diagram
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Architectuur", "End-to-end systeemoverzicht", 5)

# Draw boxes for the architecture
def arch_box(x, y, w, h, label, sub, color, label_color=WHITE):
    card(s, x, y, w, h, color, accent=GOLD)
    text_box(s, x + 0.15, y + 0.2, w - 0.3, 0.4, label, font=HEADER_FONT, size=15,
             bold=True, color=label_color, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.15, y + 0.7, w - 0.3, 0.6, sub, font=MONO_FONT, size=10,
             color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# GUI level
arch_box(0.8, 1.8, 3.2, 1.0, "GUI Bank 1", "port 8089 · public/", NAVY_LIGHT)
arch_box(9.2, 1.8, 3.2, 1.0, "GUI Bank 2", "port 8090 · public/", NAVY_LIGHT)

# API level
arch_box(0.8, 3.2, 3.2, 1.6, "Express API · Bank 1", "routes/ services/ jobs/ lib/", NAVY_LIGHT)
arch_box(9.2, 3.2, 3.2, 1.6, "Express API · Bank 2", "routes/ services/ jobs/ lib/", NAVY_LIGHT)

# DB level
arch_box(0.8, 5.2, 3.2, 1.0, "MySQL pingfin_b1", "20 accounts · 8 tables", NAVY_LIGHT)
arch_box(9.2, 5.2, 3.2, 1.0, "MySQL pingfin_b2", "20 accounts · 8 tables", NAVY_LIGHT)

# Center: Clearing Bank
arch_box(5.3, 3.4, 2.7, 1.4, "stevenop.be", "Clearing Bank · CB API v2", GOLD, NAVY)
text_box(s, 5.3, 4.95, 2.7, 0.4, "51 banks registered", font=MONO_FONT, size=11,
         color=GOLD, align=PP_ALIGN.CENTER)

# Arrows between
def arrow(x1, y1, x2, y2):
    line = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = ICE_BLUE
    line.line.width = Pt(2)

arrow(2.4, 2.8, 2.4, 3.2)   # GUI1 -> API1
arrow(10.8, 2.8, 10.8, 3.2)  # GUI2 -> API2
arrow(2.4, 4.8, 2.4, 5.2)    # API1 -> DB1
arrow(10.8, 4.8, 10.8, 5.2)  # API2 -> DB2
arrow(4.0, 4.0, 5.3, 4.05)   # API1 -> CB
arrow(8.0, 4.05, 9.2, 4.0)   # CB -> API2

# Bottom note
text_box(s, 0.6, 6.5, 12.1, 0.4,
         "Beide banken: identieke codebase · enkel .env (BIC + token) verschilt · live op Railway HTTPS",
         font=BODY_FONT, size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER, bold=True)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 6 — Tech stack
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Tech Stack", "Wat we gebruikt hebben en waarom", 6)

stacks = [
    ("Backend", "Node.js 20", "Express", "async/await · poll-jobs · groot ecosysteem"),
    ("Database", "MySQL 8", "mysql2", "DECIMAL voor saldo · prepared statements"),
    ("Frontend", "Vanilla JS", "HTML/CSS3", "Geen framework · 1000+ regels custom code"),
    ("Container", "Docker", "Compose", "Reproduceerbare deployment per bank"),
    ("Hosting", "Railway", "HTTPS", "Auto-deploy uit GitHub · 24/7 uptime"),
    ("CI Tests", "Node tests", "node-fetch", "40 unit + 18 integration tests"),
]
for i, (cat, name, tool, why) in enumerate(stacks):
    col = i % 3
    row = i // 3
    cx = 0.6 + col * 4.18
    cy = 1.8 + row * 2.5
    card(s, cx, cy, 4.0, 2.2, NAVY_LIGHT, accent=GOLD)
    text_box(s, cx + 0.3, cy + 0.2, 3.6, 0.4, cat, font=BODY_FONT, size=11,
             bold=True, color=GOLD)
    text_box(s, cx + 0.3, cy + 0.55, 3.6, 0.5, name, font=HEADER_FONT, size=22,
             bold=True, color=WHITE)
    text_box(s, cx + 0.3, cy + 1.05, 3.6, 0.4, tool, font=MONO_FONT, size=13,
             color=ICE_BLUE)
    text_box(s, cx + 0.3, cy + 1.5, 3.6, 0.7, why, font=BODY_FONT, size=11,
             color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 7 — Database schema
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Database Schema", "8 tabellen, manual-conform", 7)

tables = [
    ("accounts",      "rekeningen + saldo",            "20×€5000 startsaldo"),
    ("po_new",        "te verwerken PO's",              "input queue"),
    ("po_out",        "uitgaande PO's (wij = OB)",      "status: pending/processed/failed/timeout"),
    ("po_in",         "inkomende PO's (wij = BB)",      "verwerkte berichten"),
    ("ack_in",        "ACK's terug van CB",             "OB-zijde feedback"),
    ("ack_out",       "te versturen ACK's",             "sent_to_cb flag voor retry"),
    ("transactions",  "audit trail per saldobeweging",  "signed amount + isvalid/iscomplete"),
    ("logs",          "event-log met PO-snapshot",      "5000+ events vastgelegd"),
]
for i, (name, desc, note) in enumerate(tables):
    col = i % 2
    row = i // 2
    cx = 0.6 + col * 6.3
    cy = 1.85 + row * 1.25
    card(s, cx, cy, 6.0, 1.1, NAVY_LIGHT, accent=GOLD)
    text_box(s, cx + 0.3, cy + 0.2, 2.0, 0.4, name, font=MONO_FONT, size=15,
             bold=True, color=GOLD)
    text_box(s, cx + 0.3, cy + 0.6, 5.4, 0.4, desc, font=BODY_FONT, size=12,
             color=TEXT_LIGHT)
    text_box(s, cx + 2.4, cy + 0.2, 3.4, 0.4, note, font=BODY_FONT, size=10,
             color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

text_box(s, 0.6, 7.0, 12.1, 0.3,
         "Atomic transactions: BEGIN → 4 writes → COMMIT/ROLLBACK · geen half-verwerkte staat mogelijk",
         font=BODY_FONT, size=11, color=GOLD, align=PP_ALIGN.CENTER, bold=True)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 8 — API Endpoints
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "API Endpoints", "Public (manual) + Internal (eigen design)", 8)

# Public endpoints
text_box(s, 0.6, 1.7, 6, 0.4, "PUBLIC (manual-conform)", font=HEADER_FONT, size=14,
         bold=True, color=GOLD)
public_eps = [
    ("GET",  "/api/help",       "—"),
    ("GET",  "/api/info",       "—"),
    ("GET",  "/api/accounts",   "—"),
    ("GET",  "/api/banks",      "—"),
    ("POST", "/api/po_in",      "Bearer"),
    ("POST", "/api/ack_in",     "Bearer"),
]
for i, (method, url, auth) in enumerate(public_eps):
    cy = 2.1 + i * 0.55
    card(s, 0.6, cy, 6.0, 0.5, NAVY_LIGHT)
    color = GREEN_OK if method == "GET" else GOLD
    chip(s, 0.75, cy + 0.09, method, color, NAVY)
    text_box(s, 1.55, cy + 0.04, 3.5, 0.42, url, font=MONO_FONT, size=12,
             color=ICE_BLUE, valign=MSO_ANCHOR.MIDDLE)
    if auth != "—":
        chip(s, 5.3, cy + 0.09, auth, RED_ERR, WHITE)

# Internal endpoints
text_box(s, 6.85, 1.7, 6, 0.4, "INTERNAL (eigen design)", font=HEADER_FONT, size=14,
         bold=True, color=GOLD)
internal_eps = [
    ("GET",  "/po_new/generate",  "?count=N"),
    ("POST", "/po_new/add",       "list"),
    ("POST", "/po_new/manual",    "GUI form"),
    ("GET",  "/po_new/process",   "trigger"),
    ("GET",  "/po_out · /po_in",  "read"),
    ("GET",  "/transactions",     "audit"),
    ("GET",  "/logs",             "?type= ?limit="),
    ("GET",  "/jobs/run/:name",   "manual job"),
]
for i, (method, url, sub) in enumerate(internal_eps):
    cy = 2.1 + i * 0.55
    card(s, 6.85, cy, 6.0, 0.5, NAVY_LIGHT)
    color = GREEN_OK if method == "GET" else GOLD
    chip(s, 7.0, cy + 0.09, method, color, NAVY)
    text_box(s, 7.8, cy + 0.04, 3.5, 0.42, url, font=MONO_FONT, size=11,
             color=ICE_BLUE, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, 11.2, cy + 0.04, 1.5, 0.42, sub, font=MONO_FONT, size=10,
             color=TEXT_MUTED, valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

text_box(s, 0.6, 6.95, 12.1, 0.35, "Plus 5 background jobs · poll-po-out · poll-ack-out · flush-ack-out · timeout-monitor · cb-token",
         font=BODY_FONT, size=11, color=GOLD, align=PP_ALIGN.CENTER, bold=True)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 9 — PO Flow (success path)
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Payment Order Flow", "Use case 5: alle validaties slagen", 9)

# 7 steps in a row
steps = [
    ("1", "OB validate", "BIC · IBAN · bedrag · saldo"),
    ("2", "Debit OA", "atomic TX"),
    ("3", "POST /po_in", "→ CB"),
    ("4", "CB → BB", "push of poll"),
    ("5", "BB validate + credit", "atomic TX"),
    ("6", "POST /ack_in", "BB → CB"),
    ("7", "CB → OB ACK", "status processed"),
]
for i, (num, title, sub) in enumerate(steps):
    cx = 0.4 + i * 1.85
    # Circle with number
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.3), Inches(2.0), Inches(1.0), Inches(1.0))
    fill(circ, GOLD)
    circ.shadow.inherit = False
    text_box(s, cx + 0.3, 2.0, 1.0, 1.0, num, font=HEADER_FONT, size=32,
             bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, cx, 3.15, 1.6, 0.4, title, font=HEADER_FONT, size=12,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(s, cx, 3.55, 1.6, 0.6, sub, font=BODY_FONT, size=10,
             color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    # Arrow to next
    if i < len(steps) - 1:
        line = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(cx + 1.35), Inches(2.4),
                                  Inches(0.45), Inches(0.2))
        fill(line, ICE_BLUE)
        line.shadow.inherit = False

# Code result
card(s, 0.6, 4.7, 12.1, 1.5, NAVY_LIGHT, accent=GREEN_OK)
text_box(s, 0.85, 4.85, 11, 0.4, "Resultaat", font=HEADER_FONT, size=14,
         bold=True, color=GREEN_OK)
text_box(s, 0.85, 5.25, 11.5, 0.45, "ob_code: 2000  ·  cb_code: 2000  ·  bb_code: 2000",
         font=MONO_FONT, size=14, color=ICE_BLUE)
text_box(s, 0.85, 5.7, 11.5, 0.45, "OA -€amount  ·  BA +€amount  ·  status: processed",
         font=MONO_FONT, size=14, color=GREEN_OK)

text_box(s, 0.6, 6.5, 12.1, 0.4,
         "Atomische TX op beide banken: alles slaagt of alles wordt teruggedraaid",
         font=BODY_FONT, size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER, bold=True)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 10 — Error Handling
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Error Handling", "Wanneer wordt geld gerefund?", 10)

errors = [
    ("OB-fail", "Validatie faalt vóór debit",
     "Geen TX · geen geld bewogen", GREEN_OK),
    ("CB-fail", "CB rejecteert (4002-4007)",
     "Inline refund OA · status='failed'", AMBER),
    ("BB-fail", "BB stuurt negatieve ACK",
     "Refund OA via ackInService", AMBER),
    ("Timeout", "Geen ACK binnen 1u",
     "Auto-refund via timeoutMonitor", RED_ERR),
]
for i, (name, when, action, color) in enumerate(errors):
    col = i % 2
    row = i // 2
    cx = 0.6 + col * 6.3
    cy = 1.8 + row * 2.4
    card(s, cx, cy, 6.0, 2.1, NAVY_LIGHT, accent=color)
    text_box(s, cx + 0.3, cy + 0.2, 4.5, 0.5, name, font=HEADER_FONT, size=22,
             bold=True, color=color)
    text_box(s, cx + 0.3, cy + 0.75, 5.4, 0.4, when, font=BODY_FONT, size=13,
             color=WHITE, bold=True)
    text_box(s, cx + 0.3, cy + 1.2, 5.4, 0.8, action, font=BODY_FONT, size=12,
             color=TEXT_MUTED)

text_box(s, 0.6, 6.6, 12.1, 0.4,
         "5 background jobs zorgen dat geen enkele PO blijft hangen — self-healing system",
         font=BODY_FONT, size=12, color=GOLD, align=PP_ALIGN.CENTER, bold=True)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 11 — Code highlight: Validation
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Code Highlight 1: Validation", "Guard clauses · early return · pure functions", 11)

text_box(s, 0.6, 1.8, 6, 0.4, "Aanpak: top-down checklist", font=HEADER_FONT, size=16,
         bold=True, color=GOLD)
text_box(s, 0.6, 2.2, 6, 2.4,
         "We gebruiken early-return guard clauses i.p.v. geneste if-else.\n\n"
         "Voordelen:\n"
         "  • Leest als de manual-checklist\n"
         "  • Pure functions in lib/validate.js\n"
         "  • Herbruikbaar tussen OB- en BB-zijde\n"
         "  • Eenvoudig unit-testbaar (40/40 tests groen)",
         font=BODY_FONT, size=13, color=TEXT_LIGHT)

# Code block
card(s, 6.7, 1.8, 6.0, 4.6, DARK_BG, accent=GOLD)
text_box(s, 6.85, 1.95, 5.7, 0.4, "services/poInService.js",
         font=MONO_FONT, size=11, color=GOLD, bold=True)

code = (
    "if (po.bb_id !== BIC) {\n"
    "  await persistRejection(po, 4004);\n"
    "  return { bb_code: 4004 };\n"
    "}\n\n"
    "if (!validPoIdFormat(po.po_id)) {\n"
    "  await persistRejection(po, 4006);\n"
    "  return { bb_code: 4006 };\n"
    "}\n\n"
    "if (!validIban(po.ba_id)) {\n"
    "  await persistRejection(po, 4101);\n"
    "  return { bb_code: 4101 };\n"
    "}\n\n"
    "// alle validaties OK → atomic credit ✅"
)
text_box(s, 6.85, 2.4, 5.7, 3.9, code, font=MONO_FONT, size=11, color=ICE_BLUE)

text_box(s, 0.6, 6.55, 12.1, 0.5,
         "BIC · IBAN (15-34 chars · mod-97 checksum) · po_id · bedrag · saldo · 7 checks per PO",
         font=BODY_FONT, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 12 — Code highlight: Atomic Transactions
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Code Highlight 2: Atomic TX", "Geen half-verwerkte betalingen mogelijk", 12)

card(s, 0.6, 1.8, 12.1, 4.5, DARK_BG, accent=GOLD)
text_box(s, 0.85, 1.95, 11.5, 0.4, "services/poInService.js — credit BA atomically",
         font=MONO_FONT, size=11, color=GOLD, bold=True)

code = (
    "const conn = await pool.getConnection();\n"
    "try {\n"
    "  await conn.beginTransaction();\n"
    "  await conn.query('UPDATE accounts SET balance = balance + ? WHERE id = ?', [amount, ba_id]);\n"
    "  await conn.query('INSERT INTO po_in (...)');\n"
    "  await conn.query('INSERT INTO transactions (...)');\n"
    "  await conn.query('INSERT INTO ack_out (...)');\n"
    "  await conn.commit();   //  ✅  alles gecommit als één eenheid\n"
    "} catch (err) {\n"
    "  await conn.rollback(); //  ❌  faalt iets? alles wordt teruggedraaid\n"
    "  await persistRejection(po, 4101);\n"
    "}"
)
text_box(s, 0.85, 2.4, 11.5, 3.7, code, font=MONO_FONT, size=12, color=ICE_BLUE)

# Bottom: 4 features
feats = [
    ("✅", "ACID", "atomic"),
    ("🔒", "Lock", "row-level"),
    ("⚡", "Fast", "single conn"),
    ("🛡️", "Safe", "rollback"),
]
for i, (icon, name, desc) in enumerate(feats):
    cx = 0.6 + i * 3.05
    text_box(s, cx, 6.5, 0.6, 0.5, icon, font=HEADER_FONT, size=22, color=GOLD,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, cx + 0.6, 6.5, 2.4, 0.3, name, font=HEADER_FONT, size=13, bold=True, color=WHITE)
    text_box(s, cx + 0.6, 6.8, 2.4, 0.3, desc, font=BODY_FONT, size=10, color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 13 — Background jobs
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Background Jobs", "Self-healing system · 5 processen", 13)

jobs = [
    ("pollPoOut",       "30s",  "Haalt PO's voor onze BIC uit CB.po_out"),
    ("pollAckOut",      "30s",  "Haalt ACK's voor onze verstuurde PO's uit CB.ack_out"),
    ("flushAckOut",     "60s",  "Retry voor ACK's die niet meteen door CB werden geaccepteerd"),
    ("timeoutMonitor",  "5min", "Outstanding po_out > 1u → refund OA + status='timeout'"),
    ("cbTokenRefresh",  "3.5h", "Vernieuwt Bearer-token vóór TTL van 4u verloopt"),
]
for i, (name, freq, desc) in enumerate(jobs):
    cy = 1.8 + i * 1.0
    card(s, 0.6, cy, 12.1, 0.85, NAVY_LIGHT, accent=GOLD)
    text_box(s, 0.85, cy + 0.13, 3.5, 0.5, name, font=MONO_FONT, size=16,
             bold=True, color=GOLD, valign=MSO_ANCHOR.MIDDLE)
    chip(s, 4.4, cy + 0.27, freq, ICE_BLUE, NAVY)
    text_box(s, 5.6, cy + 0.13, 7.0, 0.5, desc, font=BODY_FONT, size=13,
             color=TEXT_LIGHT, valign=MSO_ANCHOR.MIDDLE)

text_box(s, 0.6, 6.95, 12.1, 0.4,
         "Resultaat: als CB even down is, herstelt het systeem zichzelf zonder data te verliezen",
         font=BODY_FONT, size=12, color=GOLD, align=PP_ALIGN.CENTER, bold=True)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 14 — GUI features
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "GUI Features", "Gebruiksvriendelijk + realtime", 14)

features = [
    ("📊 Dashboard",       "6 stats + Quick Actions + bank info"),
    ("💳 Accounts",         "Live saldo's met euro-formatting"),
    ("➕ Manuele PO",       "Slimme dropdowns: OA met saldo · BB uit CB.banks · BA datalist"),
    ("📤📥 PO_OUT/PO_IN",   "Realtime tabellen, kleurgecodeerde badges per code"),
    ("✅📨 ACK tabs",       "Inkomende + verstuurde ACK's met datum-locale nl-BE"),
    ("💱 Transacties",      "Audit-trail met +/- bedragen + valid/complete vlaggen"),
    ("📜 Logs",             "Filter op event-type · 17 verschillende types"),
    ("🏛️ Banks",            "51 banken uit CB-cache, eigen BIC gehighlight"),
]
for i, (icon_title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    cx = 0.6 + col * 6.3
    cy = 1.8 + row * 1.25
    card(s, cx, cy, 6.0, 1.1, NAVY_LIGHT, accent=GOLD)
    text_box(s, cx + 0.3, cy + 0.18, 5.5, 0.45, icon_title, font=HEADER_FONT, size=15,
             bold=True, color=WHITE)
    text_box(s, cx + 0.3, cy + 0.6, 5.5, 0.5, desc, font=BODY_FONT, size=11,
             color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 15 — Live notifications
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Live Notifications", "Auto-poll elke 10s · diff-detection · toast popups", 15)

# Left column: how it works
text_box(s, 0.6, 1.8, 6.0, 0.4, "Hoe werkt het?", font=HEADER_FONT, size=16,
         bold=True, color=GOLD)
text_box(s, 0.6, 2.25, 6.0, 4.0,
         "Elke 10 seconden:\n\n"
         "1.  GUI haalt po_in / po_out / ack_in / ack_out / accounts op\n\n"
         "2.  Vergelijkt met snapshot van vorige poll\n\n"
         "3.  Toont toast popup voor élke nieuwe rij of statuswijziging\n\n"
         "4.  Pulsende dot op nav-knop voor ongelezen events\n\n"
         "5.  Live indicator linksonder: groen=OK, rood=API down",
         font=BODY_FONT, size=13, color=TEXT_LIGHT)

# Right: example toasts
text_box(s, 7.0, 1.8, 5.7, 0.4, "Voorbeeld events", font=HEADER_FONT, size=16,
         bold=True, color=GOLD)
toasts = [
    ("📥", "Nieuwe inkomende PO verwerkt", "CEKVBE88_a1b2 · €50",   GREEN_OK),
    ("✅", "ACK ontvangen",                 "CEKVBE88_a1b2 · 2000",  GREEN_OK),
    ("✕",  "Negatieve ACK (4004)",          "CEKVBE88_x9y2",          RED_ERR),
    ("💰", "Saldo gewijzigd",               "BE13... €5000 → €4950", AMBER),
    ("⏰", "PO getimeout (1u)",             "CEKVBE88_z3w4",          RED_ERR),
]
for i, (icon, title, sub, color) in enumerate(toasts):
    cy = 2.3 + i * 0.85
    card(s, 7.0, cy, 5.7, 0.7, NAVY_LIGHT, accent=color)
    text_box(s, 7.15, cy + 0.12, 0.6, 0.5, icon, font=HEADER_FONT, size=18,
             color=color, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, 7.75, cy + 0.07, 4.8, 0.35, title, font=BODY_FONT, size=12,
             bold=True, color=WHITE)
    text_box(s, 7.75, cy + 0.4, 4.8, 0.3, sub, font=MONO_FONT, size=10,
             color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 16 — Demo: stap 1 manual PO
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Live Demo · Stap 1", "Manuele PO via GUI · Bank 1 → Bank 2", 16)

# Left: action steps
text_box(s, 0.6, 1.8, 5.5, 0.4, "Wat we doen", font=HEADER_FONT, size=16,
         bold=True, color=GOLD)
steps = [
    "Open localhost:8089 (Bank 1)",
    "Klik tab '➕ PO Aanmaken'",
    "Open 'Manuele PO' details",
    "OA dropdown → kies eigen rekening (saldo zichtbaar)",
    "BB dropdown → kies HOMNBEB1 (Bank 2)",
    "BA IBAN → typ ontvanger",
    "Bedrag → €50",
    "Klik 'Verstuur manuele PO'",
]
for i, step in enumerate(steps):
    cy = 2.3 + i * 0.42
    # Number circle
    n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(cy + 0.05), Inches(0.3), Inches(0.3))
    fill(n, GOLD)
    n.shadow.inherit = False
    text_box(s, 0.6, cy + 0.05, 0.3, 0.3, str(i+1), font=HEADER_FONT, size=11,
             bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, 1.05, cy + 0.05, 5.0, 0.3, step, font=BODY_FONT, size=12,
             color=TEXT_LIGHT, valign=MSO_ANCHOR.MIDDLE)

# Right: screenshot placeholder
card(s, 6.5, 1.8, 6.3, 5.0, DARK_BG, accent=GOLD)
text_box(s, 6.7, 2.0, 6.0, 0.4, "[ SCREENSHOT 1: Manuele PO formulier ]",
         font=HEADER_FONT, size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
text_box(s, 6.7, 2.5, 6.0, 4.0,
         "Maak screenshot van:\n"
         "→ GUI tab 'PO Aanmaken'\n"
         "→ Manuele PO details geopend\n"
         "→ Dropdowns met eigen accounts + BB-banken zichtbaar\n"
         "→ Toast 'Manuele PO aangemaakt' rechtsonder\n\n\n"
         "Vervang deze tekst met je screenshot",
         font=BODY_FONT, size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 17 — Demo: stap 2 process + ACK
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Live Demo · Stap 2", "Verwerk + ontvang ACK", 17)

# Top: timeline
text_box(s, 0.6, 1.7, 12.1, 0.4, "Wat gebeurt er na 'Verwerk PO_NEW'?",
         font=HEADER_FONT, size=15, bold=True, color=GOLD)

timeline = [
    ("0s",   "Bank 1: PO_OUT krijgt rij · saldo daalt €50 · debit transactions"),
    ("~30s", "Bank 2 BB-poller: detecteert PO · processPoIn · credit BA · queue ack_out"),
    ("~30s", "Bank 2: 🟢 Toast 'Nieuwe inkomende PO verwerkt'"),
    ("~30s", "Bank 2: directe POST /ack_in naar CB"),
    ("~60s", "Bank 1 OB-poller: pickt ACK op uit CB.ack_out"),
    ("~60s", "Bank 1: 🟢 Toast 'ACK ontvangen' · status → processed"),
]
for i, (t, what) in enumerate(timeline):
    cy = 2.3 + i * 0.55
    chip(s, 0.6, cy + 0.07, t, GOLD, NAVY)
    text_box(s, 1.85, cy + 0.05, 11, 0.45, what, font=BODY_FONT, size=12,
             color=TEXT_LIGHT, valign=MSO_ANCHOR.MIDDLE)

# Bottom: screenshot placeholder
card(s, 0.6, 5.7, 12.1, 1.4, DARK_BG, accent=GOLD)
text_box(s, 0.85, 5.85, 11.5, 0.35, "[ SCREENSHOT 2: PO_OUT tab met rij status='processed' + Toast 'ACK ontvangen' ]",
         font=HEADER_FONT, size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
text_box(s, 0.85, 6.2, 11.5, 0.7,
         "Maak screenshot na ~60s van GUI Bank 1: PO_OUT met cb_code 2000 + bb_code 2000 + groene status badge",
         font=BODY_FONT, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 18 — Demo: error scenarios
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Live Demo · Stap 3", "Error scenarios — élke foutcode getriggerd", 18)

# Code block
card(s, 0.6, 1.8, 12.1, 2.7, DARK_BG, accent=RED_ERR)
text_box(s, 0.85, 1.95, 11.5, 0.4, "Test 4002 — bedrag te hoog",
         font=MONO_FONT, size=11, bold=True, color=GOLD)
text_box(s, 0.85, 2.4, 11.5, 0.4,
         "curl -X POST http://localhost:8089/api/po_new/manual \\",
         font=MONO_FONT, size=11, color=ICE_BLUE)
text_box(s, 0.85, 2.75, 11.5, 0.4,
         "  -d '{\"oa_id\":\"BE13...\",\"ba_id\":\"BE99...\",\"bb_id\":\"HOMNBEB1\",\"po_amount\":600}'",
         font=MONO_FONT, size=11, color=ICE_BLUE)
text_box(s, 0.85, 3.3, 11.5, 0.4, "Response:",
         font=MONO_FONT, size=11, bold=True, color=GOLD)
text_box(s, 0.85, 3.7, 11.5, 0.6,
         '{ "ok": false, "code": 4002, "message": "Ongeldig bedrag" }',
         font=MONO_FONT, size=12, color=RED_ERR)

# Bottom: codes tested
text_box(s, 0.6, 4.7, 12.1, 0.4, "Alle 10 foutcodes worden in tests/error-pos.test.js geverifieerd",
         font=HEADER_FONT, size=14, bold=True, color=GOLD)
codes = [
    "4002 bedrag>500", "4003 bedrag<0", "4004 BB unknown", "4005 duplicate",
    "4006 OB mismatch", "4101 ACCOUNT_UNKNOWN", "4102 saldo<0", "401 no Bearer", "401 wrong Bearer"
]
for i, code in enumerate(codes):
    col = i % 3
    row = i // 3
    cx = 0.6 + col * 4.18
    cy = 5.2 + row * 0.55
    chip(s, cx, cy, code, RED_ERR, WHITE)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 19 — Test suite
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Test Suite", "58 geautomatiseerde tests · alles groen", 19)

# Big stats
big_number(s, 0.6, 1.7, "40", "Unit tests", GOLD, 90, 16)
big_number(s, 4.0, 1.7, "18", "Integration tests", GOLD, 90, 16)
big_number(s, 7.4, 1.7, "100%", "Pass rate", GREEN_OK, 70, 16)
big_number(s, 11.0, 1.7, "10", "Foutcodes gedekt", ICE_BLUE, 90, 16)

# Tests overview cards
card(s, 0.6, 4.5, 6.0, 2.4, NAVY_LIGHT, accent=GOLD)
text_box(s, 0.85, 4.65, 5.7, 0.4, "Unit tests (validate.test.js)",
         font=HEADER_FONT, size=15, bold=True, color=GOLD)
text_box(s, 0.85, 5.05, 5.7, 1.9,
         "•  BIC validatie (8/11 chars · case-insensitive)\n"
         "•  IBAN validatie (15-34 chars · mod-97 checksum)\n"
         "•  po_id format (BIC prefix · max 50 chars)\n"
         "•  Bedrag (>0 · ≤500 · max 2 decimalen)\n"
         "•  Datum format YYYY-MM-DD HH:MM:SS\n"
         "•  Alle 10 error codes correct\n\n"
         "$  npm test  →  40 / 0",
         font=BODY_FONT, size=11, color=TEXT_LIGHT)

card(s, 6.7, 4.5, 6.0, 2.4, NAVY_LIGHT, accent=GREEN_OK)
text_box(s, 6.95, 4.65, 5.7, 0.4, "Integration tests (error-pos.test.js)",
         font=HEADER_FONT, size=15, bold=True, color=GREEN_OK)
text_box(s, 6.95, 5.05, 5.7, 1.9,
         "•  POST manual PO met opzettelijke fout\n"
         "•  Verifieer juiste error-code in response\n"
         "•  Verifieer log-event geschreven in /api/logs\n"
         "•  10 scenarios + 8 assertions per scenario\n"
         "•  Bearer-auth: 401 zonder/met verkeerd token\n"
         "•  Happy-path baseline\n\n"
         "$  npm run test:errors  →  18 / 0",
         font=BODY_FONT, size=11, color=TEXT_LIGHT)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 20 — Foutcodes table
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Foutcodes", "10 codes — manual-conform + eigen extensies", 20)

codes_data = [
    ("2000", "OK",                    "succesvolle verwerking",                    GREEN_OK),
    ("4001", "INTERNAL_TX",           "interne PO foutief naar CB gestuurd",       AMBER),
    ("4002", "AMOUNT_EXCEEDED",       "bedrag > €500",                             RED_ERR),
    ("4003", "AMOUNT_INVALID",        "bedrag ≤ 0 of NaN",                         RED_ERR),
    ("4004", "BB_UNKNOWN",            "ontvangende BIC niet bekend",               RED_ERR),
    ("4005", "DUPLICATE_PO",          "po_id reeds verwerkt",                      AMBER),
    ("4006", "OB_MISMATCH",           "po_id-prefix klopt niet",                   AMBER),
    ("4007", "DUP_IN_BATCH",          "dezelfde po_id 2× in 1 batch",              AMBER),
    ("4101", "ACCOUNT_UNKNOWN",       "OA/BA bestaat niet of ongeldige IBAN",      RED_ERR),
    ("4102", "INSUFFICIENT_BALANCE",  "OA-saldo < bedrag",                         RED_ERR),
]
for i, (code, name, when, color) in enumerate(codes_data):
    col = i % 2
    row = i // 2
    cx = 0.6 + col * 6.3
    cy = 1.75 + row * 1.0
    card(s, cx, cy, 6.0, 0.85, NAVY_LIGHT, accent=color)
    text_box(s, cx + 0.2, cy + 0.15, 0.9, 0.5, code, font=MONO_FONT, size=20,
             bold=True, color=color, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, cx + 1.25, cy + 0.1, 4.6, 0.3, name, font=MONO_FONT, size=11,
             bold=True, color=ICE_BLUE)
    text_box(s, cx + 1.25, cy + 0.42, 4.6, 0.4, when, font=BODY_FONT, size=10,
             color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 21 — Security
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Security", "8 verdedigingslagen · OWASP-conform", 21)

security = [
    ("🔑", "Bearer Auth",      "Inkomend POST → 401 zonder token"),
    ("🔄", "CB-token rotatie", "Auto-refresh elke 3.5h (TTL 4h)"),
    ("💉", "SQL injection",    "mysql2 prepared statements · 0 string-concat"),
    ("⚡", "XSS prevention",    "escapeHtml() op alle user-data in GUI"),
    ("🛡️", "CSP + headers",     "X-Frame-Options · script-src 'self'"),
    ("🚦", "Rate limit",       "60 POST/min/IP · 429 op overflow"),
    ("🙈", "Hidden links",     "Trello/GitHub alleen op localhost"),
    ("🤐", "Secrets",          "INCOMING_TOKEN in .env (git-ignored)"),
]
for i, (icon, title, desc) in enumerate(security):
    col = i % 2
    row = i // 2
    cx = 0.6 + col * 6.3
    cy = 1.8 + row * 1.25
    card(s, cx, cy, 6.0, 1.1, NAVY_LIGHT, accent=GOLD)
    text_box(s, cx + 0.2, cy + 0.18, 0.7, 0.7, icon, font=HEADER_FONT, size=24,
             color=GOLD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, cx + 0.95, cy + 0.15, 4.9, 0.4, title, font=HEADER_FONT, size=13,
             bold=True, color=WHITE)
    text_box(s, cx + 0.95, cy + 0.55, 4.9, 0.5, desc, font=BODY_FONT, size=10,
             color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 22 — Difficulties
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Difficulties", "6 grote bugs · allemaal gefixt op dag 3", 22)

bugs = [
    ("OB krijgt geen ACK",        "Silent failures in 4 reject-paden",
     "persistRejection() helper"),
    ("4101-storm",                 "IBAN regex was BE-only (16 chars)",
     "Regex naar 15-34 chars internationaal"),
    ("4004 zelf veroorzaakt",      "Hard-coded BIC fallback in generator",
     "503 fail-fast bij CB-onbereik"),
    ("CB-rejectie 1u limbo",       "po_out bleef pending bij 4xxx",
     "Inline refund in poProcessor"),
    ("flushAckOut dropte rijen",   "LEFT JOIN faalde bij ontbrekende po_in",
     "Fallback naar logs-snapshot"),
    ("BIC case-sensitive",         "Andere teams stuurden lowercase",
     "toUpperCase() in alle vergelijkingen"),
]
for i, (title, problem, fix) in enumerate(bugs):
    col = i % 2
    row = i // 2
    cx = 0.6 + col * 6.3
    cy = 1.8 + row * 1.7
    card(s, cx, cy, 6.0, 1.55, NAVY_LIGHT, accent=RED_ERR)
    text_box(s, cx + 0.25, cy + 0.15, 5.5, 0.4, "🐛  " + title, font=HEADER_FONT, size=13,
             bold=True, color=RED_ERR)
    text_box(s, cx + 0.25, cy + 0.55, 5.5, 0.4, problem, font=BODY_FONT, size=11,
             color=TEXT_MUTED)
    text_box(s, cx + 0.25, cy + 0.95, 5.5, 0.5, "✅  " + fix, font=BODY_FONT, size=11,
             color=GREEN_OK, bold=True)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 23 — Lessons learned
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Lessons Learned", "Wat we de volgende keer anders doen", 23)

lessons = [
    ("✅", "Wat ging goed",
     "Vroege CB-API simulatie (dag 1)\n"
     "Atomische TX vanaf dag 2\n"
     "Background jobs voor veerkracht\n"
     "Test suite voor regressie-detectie\n"
     "Live notifications: 10× sneller debuggen"),
    ("⚠️", "Wat ging minder",
     "OneDrive + git = sync-conflicten\n"
     "Hard-coded fallbacks zijn anti-pattern\n"
     "Silent error swallowing verbergde bugs\n"
     "IBAN regex te strikt (BE-only)\n"
     "Bug-fixes liepen door tot dag 4"),
    ("🚀", "Volgende keer",
     "CI/CD met automated tests vanaf dag 1\n"
     "Vroege deployment (dag 2)\n"
     "Geen project in OneDrive — alleen git\n"
     "Pair-programming voor cross-functional kennis\n"
     "API-versionering vanaf het begin"),
]
for i, (icon, title, items) in enumerate(lessons):
    cx = 0.6 + i * 4.18
    card(s, cx, 1.8, 4.0, 5.0, NAVY_LIGHT, accent=GOLD)
    text_box(s, cx + 0.2, 1.95, 0.7, 0.7, icon, font=HEADER_FONT, size=28,
             color=GOLD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, cx + 1.0, 2.0, 2.9, 0.5, title, font=HEADER_FONT, size=16,
             bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, cx + 0.2, 2.85, 3.6, 3.9, items, font=BODY_FONT, size=11,
             color=TEXT_LIGHT)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 24 — Eindstand
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
page_header(s, "Eindstand · Wat is gerealiseerd?", "Alle deliverables van dag 1-4", 24)

# Big stats top
big_number(s, 0.6, 1.7, "1066+", "ACK's verwerkt",  GREEN_OK, 70, 14)
big_number(s, 3.7, 1.7, "51",    "Banks bij CB",    ICE_BLUE, 70, 14)
big_number(s, 6.8, 1.7, "58",    "Tests groen",     GREEN_OK, 70, 14)
big_number(s, 9.9, 1.7, "24/7",  "Online",          GOLD,     60, 14)

# Checklist
done = [
    "✅ Beide banken volledig functioneel — verzenden + ontvangen",
    "✅ Alle 5 use cases uit de manual gedekt",
    "✅ 5 background jobs draaien (poll-po-out · poll-ack-out · flush · timeout · token)",
    "✅ 40 unit + 18 integration tests, allemaal groen",
    "✅ GUI met live updates, toast notifications, slimme dropdowns",
    "✅ Beide banken live op Railway met HTTPS",
    "✅ Per-bank Docker-folder structuur (coach-eis)",
    "✅ Security: 8 verdedigingslagen + audit document",
    "✅ 10 documentatie files (4 dag-rapporten + verslag + tests + security)",
]
for i, item in enumerate(done):
    cy = 4.0 + i * 0.36
    text_box(s, 0.6, cy, 12.1, 0.32, item, font=BODY_FONT, size=12,
             color=TEXT_LIGHT, valign=MSO_ANCHOR.MIDDLE)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 25 — Q&A / Thanks
# ──────────────────────────────────────────────────────────────────────
s = add_blank(prs)
bg(s, DARK_BG)

# Decorative gold band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SH)
fill(band, GOLD)
band.shadow.inherit = False

text_box(s, 1.0, 1.8, 11, 1.5, "Bedankt!", font=HEADER_FONT, size=80, bold=True, color=GOLD)
text_box(s, 1.0, 3.2, 11, 0.6, "Vragen?", font=HEADER_FONT, size=36, color=WHITE)

# Live links cards
text_box(s, 1.0, 4.2, 11, 0.4, "Probeer het zelf:", font=BODY_FONT, size=14,
         bold=True, color=TEXT_MUTED)

card(s, 1.0, 4.7, 5.5, 0.9, NAVY_LIGHT, accent=GOLD)
text_box(s, 1.2, 4.85, 5.1, 0.3, "Bank 1 · CEKVBE88", font=HEADER_FONT, size=12,
         bold=True, color=GOLD)
text_box(s, 1.2, 5.15, 5.1, 0.4, "pingfin-team20-production.up.railway.app",
         font=MONO_FONT, size=11, color=ICE_BLUE)

card(s, 6.8, 4.7, 5.5, 0.9, NAVY_LIGHT, accent=GOLD)
text_box(s, 7.0, 4.85, 5.1, 0.3, "Bank 2 · HOMNBEB1", font=HEADER_FONT, size=12,
         bold=True, color=GOLD)
text_box(s, 7.0, 5.15, 5.1, 0.4, "pingfin-team20-bank2-production.up.railway.app",
         font=MONO_FONT, size=11, color=ICE_BLUE)

text_box(s, 1.0, 6.4, 11, 0.4,
         "Salaheddine Sennouni  ·  Abdallah Azouagh  ·  Ayoub Abdeddoun  ·  Marwan Saidi",
         font=BODY_FONT, size=14, color=TEXT_MUTED)
text_box(s, 1.0, 6.85, 11, 0.3,
         "Workshop Software Engineering 2026  ·  Odisee  ·  PingFin Team 20",
         font=BODY_FONT, size=11, color=TEXT_MUTED)


# ──────────────────────────────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────────────────────────────
out = "docs/PingFin-Team20-Presentation.pptx"
prs.save(out)
print(f"OK Presentation saved: {out} ({len(prs.slides)} slides)")
